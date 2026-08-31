"""PPTX file ingestor."""

from __future__ import annotations

from pathlib import Path

from app.domain.documents import DocumentMetadata, SourceDocument
from app.infrastructure.ingestion.base import (
    BaseIngestor,
    IngestionError,
    SourceReference,
    require_path_source,
)
from app.infrastructure.ingestion.utils import file_timestamp


class PptxIngestor(BaseIngestor):
    """Read PPTX files into normalized source documents."""

    source_type = "pptx"
    supported_suffixes = (".pptx", ".ppt", ".odp")

    def ingest(self, source: SourceReference) -> SourceDocument:
        source_path = require_path_source(source, ingestor_name="PPTX ingestor")
        text = self._extract_text(source_path)
        resolved_path = source_path.resolve()
        suffix = source_path.suffix.lower()
        if suffix == ".odp":
            mime = "application/vnd.oasis.opendocument.presentation"
        elif suffix == ".ppt":
            mime = "application/vnd.ms-powerpoint"
        else:
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        return SourceDocument(
            source=str(resolved_path),
            source_path=resolved_path,
            source_type=self.source_type,
            filename=source_path.name,
            text=text,
            metadata=DocumentMetadata(
                title=source_path.stem,
                modified_at=file_timestamp(source_path),
                mime_type=mime,
            ),
        )

    def _extract_text(self, path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix == ".odp":
            return self._extract_odp(path)
        if suffix == ".ppt":
            return self._extract_ppt(path)
        return self._extract_pptx(path)

    def _extract_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore[import-untyped]

            prs = Presentation(str(path))
            texts: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            raise IngestionError(
                "python-pptx is required for PPTX ingestion. "
                "Install with: pip install python-pptx"
            ) from None
        except Exception as exc:
            raise IngestionError(f"Unable to read PPTX file '{path}'.") from exc

    def _extract_odp(self, path: Path) -> str:
        try:
            from odf.opendocument import load  # type: ignore[import-untyped]
            from odf import draw, text as odf_text, teletype  # type: ignore[import-untyped]

            doc = load(str(path))
            texts: list[str] = []
            for page in doc.getElementsByType(draw.Page):
                for frame in page.getElementsByType(draw.Frame):
                    for tb in frame.getElementsByType(draw.TextBox):
                        for p in tb.getElementsByType(odf_text.P):
                            t = teletype.extractText(p)
                            if t.strip():
                                texts.append(t)
            return "\n".join(texts)
        except ImportError:
            raise IngestionError(
                "odfpy is required for ODP ingestion. Install with: pip install odfpy"
            ) from None
        except IngestionError:
            raise
        except Exception as exc:
            raise IngestionError(f"Unable to read ODP file '{path}'.") from exc

    def _extract_ppt(self, path: Path) -> str:
        try:
            import struct

            import olefile  # type: ignore[import-untyped]

            if not olefile.isOleFile(str(path)):
                raise IngestionError(f"File '{path}' is not a valid OLE2 file.")

            ole = olefile.OleFileIO(str(path))
            texts: list[str] = []
            if ole.exists("PowerPoint Document"):
                data = ole.openstream("PowerPoint Document").read()
                offset = 0
                while offset < len(data) - 8:
                    try:
                        rec_type, rec_len = struct.unpack_from("<II", data, offset)
                    except struct.error:
                        break
                    if rec_len < 0 or offset + 8 + rec_len > len(data):
                        offset += 1
                        continue
                    rec_data = data[offset + 8 : offset + 8 + rec_len]
                    if rec_type == 0x0FA8:
                        t = rec_data.decode("latin-1", errors="replace").strip()
                        if t:
                            texts.append(t)
                    elif rec_type == 0x0FA9:
                        t = rec_data.decode("utf-16-le", errors="replace").strip()
                        if t:
                            texts.append(t)
                    offset += 8 + rec_len
            ole.close()
            return "\n".join(texts)
        except ImportError:
            raise IngestionError(
                "olefile is required for PPT ingestion. Install with: pip install olefile"
            ) from None
        except IngestionError:
            raise
        except Exception as exc:
            raise IngestionError(f"Unable to read PPT file '{path}'.") from exc
